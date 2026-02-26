from pptx import Presentation
from pptx.util import Inches, Pt

try:
    prs = Presentation()

    # Slide 1 - Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Agno Agent 簡介"
    subtitle.text = "快速、多代理框架 — 摘要與資源"

    # Slide 2 - What is Agno?
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = "什麼是 Agno?"
    body.text = "- 輕量且快速的多代理 (multi-agent) 框架"
    p = body.add_paragraph()
    p.text = "- 支援工具整合、記憶體、知識管理與推理"
    p.level = 1
    p = body.add_paragraph()
    p.text = "- 適合開發可部署於生產的 agent 應用"
    p.level = 1

    # Slide 3 - 核心概念
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = "核心概念"
    body.text = "- 宣告式 agent 組裝: 以 Python 配置 model、memory、tools"
    p = body.add_paragraph()
    p.text = "- Tools as first-class citizens: 工具呼叫與外部系統整合"
    p = body.add_paragraph()
    p.text = "- 可檢視的推理流程與除錯追蹤"

    # Slide 4 - 架構與元件
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = "架構與元件"
    body.text = "- Models: 支援多種 LLM 後端 (例如 OpenAI, Anthropic 等)"
    p = body.add_paragraph()
    p.text = "- Memory: 短期會話記憶 & 長期持久化存儲"
    p = body.add_paragraph()
    p.text = "- Knowledge/RAG: 向量資料庫與檢索增強生成"
    p = body.add_paragraph()
    p.text = "- Tools & Integrations: 內建 100+ 工具套件"

    # Slide 5 - 常見應用與範例
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = "應用場景與範例"
    body.text = "- 研究助理: 多步驟檢索、整理與報告生成"
    p = body.add_paragraph()
    p.text = "- 工作流程自動化: 串接專責 agents 形成 pipeline"
    p = body.add_paragraph()
    p.text = "- 產品化部署: 可在生產環境高速執行與監控"

    # Slide 6 - 參考資源
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1].text_frame
    title.text = "參考資源"
    body.text = "- 官方 GitHub: https://github.com/agno-agi/agno"
    p = body.add_paragraph()
    p.text = "- 文章: DigitalOcean 講解 Agno 框架"
    p = body.add_paragraph()
    p.text = "- 官方文件: https://docs.agno.com (若可用)"

    output_path = "agno_agent_overview.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(output_path)

except Exception as e:
    print("Error while creating presentation:", str(e))
    raise

print('Done')