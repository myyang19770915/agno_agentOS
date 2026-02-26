from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
from PIL import Image, ImageDraw, ImageFont
import datetime

try:
    # Parameters
    date_str = '2026-02-23'  # as requested (today)
    author = 'Author: AI Team'

    # Paths
    os.makedirs('downloads', exist_ok=True)
    os.makedirs('outputs/images', exist_ok=True)
    pptx_path = 'downloads/agent_architecture_onpage_report.pptx'
    png_path = 'outputs/images/agent_architecture_onpage_report.png'

    # Create presentation with 16:9 size
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # Colors
    deep_blue = RGBColor(28,40,51)   # #1C2833
    slate_gray = RGBColor(46,64,83)  # #2E4053
    light_blue = RGBColor(94,168,167) # #5EA8A7 (used as accent)
    orange = RGBColor(249,109,0)     # #F96D00
    text_black = RGBColor(0,0,0)

    # Fonts
    title_font = 'Arial'
    body_font = 'Arial'

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # HEADER: Title and subtitle
    left_margin = Inches(0.4)
    right_margin = Inches(0.4)

    # Title box
    title_box = slide.shapes.add_textbox(left_margin, Inches(0.25), prs.slide_width - left_margin - right_margin, Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = 'Agent 架構'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.name = title_font
    p.font.color.rgb = deep_blue

    # Subtitle / meta
    subtitle_box = slide.shapes.add_textbox(left_margin, Inches(1.05), prs.slide_width - left_margin - right_margin, Inches(0.6))
    st = subtitle_box.text_frame
    st.text = f'On-page Report — 簡要概覽    {date_str}    {author}'
    for paragraph in st.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.name = body_font
        paragraph.font.color.rgb = slate_gray

    # LEFT: Architecture diagram area
    diagram_left = Inches(0.5)
    diagram_top = Inches(1.9)
    diagram_width = Inches(7.5)
    diagram_height = Inches(4.5)

    # Draw modules as rounded rectangles horizontally with arrows
    module_w = Inches(1.6)
    module_h = Inches(0.7)
    spacing = Inches(0.25)
    start_x = diagram_left + Inches(0.1)
    start_y = diagram_top + Inches(0.5)

    modules = [
        ('感知\n(Perception)', '輸入感測/資料擷取，解析原始輸入'),
        ('理解/推理\n(Reasoning)', '訊息整合、策略決策與規則推理'),
        ('記憶\n(Memory)', '長短期知識庫、情境與會話歷史'),
        ('學習\n(Learning)', '模型更新、離線/線上訓練迴路'),
        ('執行\n(Action)', '發送指令、呼叫 API 或控制致動器'),
    ]

    shape_centers = []
    cur_x = start_x
    for i,(title, desc) in enumerate(modules):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cur_x, start_y, module_w, module_h)
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor(236,243,249) if i%2==0 else RGBColor(230,240,255)
        shp.line.color.rgb = slate_gray
        # Title inside shape
        tf = shp.text_frame
        tf.text = title
        for p in tf.paragraphs:
            p.font.size = Pt(11)
            p.font.name = body_font
            p.font.bold = True
            p.font.color.rgb = deep_blue
        shape_centers.append((cur_x + module_w/2, start_y + module_h/2))
        cur_x += module_w + spacing

    # Add bidirectional arrow between Memory and Learning
    # Using simple connector arrow shapes: add_shape RIGHT_ARROW between them
    # For simplicity, add small arrows shapes between modules
    arrow_w = Inches(0.4)
    arrow_h = Inches(0.18)
    # arrows between modules
    arrow_x = start_x + module_w + spacing/2
    arrow_y = start_y + module_h/2 - arrow_h/2
    for i in range(len(modules)-1):
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, start_x + (module_w+spacing)*i + module_w + spacing/6, arrow_y, module_w/6, arrow_h)
        arr.fill.solid()
        arr.fill.fore_color.rgb = slate_gray
        arr.line.fill.background()

    # For Memory<->Learning, add double arrow
    # Add a small double-headed arrow shape above the memory-learning gap
    mem_idx = 2
    learn_idx = 3
    x_between = start_x + (module_w+spacing) * mem_idx + module_w + spacing/6
    dbl = slide.shapes.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, x_between, start_y - Inches(0.25), module_w/6, Inches(0.12))
    dbl.fill.solid()
    dbl.fill.fore_color.rgb = light_blue
    dbl.line.fill.background()

    # Add descriptions under the diagram
    desc_top = start_y + module_h + Inches(0.25)
    desc_box = slide.shapes.add_textbox(diagram_left, desc_top, diagram_width, Inches(1.6))
    d_tf = desc_box.text_frame
    d_tf.word_wrap = True
    # 利用第一個預設空 paragraph，避免 del（paragraphs 是 tuple 不可刪）
    for i, (title, desc) in enumerate(modules):
        p = d_tf.paragraphs[0] if i == 0 else d_tf.add_paragraph()
        p.text = f"{title.split(chr(10))[0]}：{desc}"
        p.level = 0
        p.font.size = Pt(9)
        p.font.name = body_font
        p.font.color.rgb = slate_gray

    # Add external interfaces: Users/Integrations and Environment
    ext_left = diagram_left - Inches(0.2)
    # Users box left side
    users_box = slide.shapes.add_textbox(ext_left, start_y - Inches(0.5), Inches(1.3), Inches(0.6))
    tb = users_box.text_frame
    tb.text = 'Users / Integrations\n(使用者 / 系統整合)'
    for p in tb.paragraphs:
        p.font.size = Pt(9)
        p.font.name = body_font
        p.font.color.rgb = dark_color = slate_gray

    # Environment box right side
    env_box = slide.shapes.add_textbox(start_x + diagram_width - Inches(0.9), start_y + module_h + Inches(0.6), Inches(1.6), Inches(0.6))
    etf = env_box.text_frame
    etf.text = 'Environment\n(執行環境 / 感測場域)'
    for p in etf.paragraphs:
        p.font.size = Pt(9)
        p.font.name = body_font
        p.font.color.rgb = slate_gray

    # RIGHT: Key points bullet list (最多6點)
    right_left = Inches(8.6)
    right_top = Inches(1.9)
    right_w = prs.slide_width - right_left - right_margin
    right_h = Inches(4.5)

    points = [
        ("模組職責明確", "各模組（感知/理解/記憶/學習/執行）單責任、易於測試與替換"),
        ("資料流設計", "採用事件/訊息驅動，明確輸入/輸出介面與資料契約"),
        ("可擴充性", "模組化與 API 化，支援水平擴展與外掛式擴充"),
        ("安全與隱私", "資料分類、存取控制、加密與審計紀錄為設計基礎"),
        ("監控與回饋", "日誌、指標、延遲與效果監控，並提供線上回饋迴路"),
        ("部署考量", "容器化、資源隔離、模型版本管理與 A/B 測試策略")
    ]

    bullet_box = slide.shapes.add_textbox(right_left, right_top, right_w, right_h)
    btf = bullet_box.text_frame
    btf.word_wrap = True
    # Add each point as bold short line + small description
    # 利用第一個預設空 paragraph，避免 del（paragraphs 是 tuple 不可刪）
    first = True
    for title, detail in points:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.text = title
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.name = body_font
        p.font.color.rgb = deep_blue
        p.level = 0
        p2 = btf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(10)
        p2.font.name = body_font
        p2.level = 1
        p2.font.color.rgb = slate_gray

    # Bottom: conclusion and contact
    conclusion_box = slide.shapes.add_textbox(left_margin, prs.slide_height - Inches(0.8), prs.slide_width - left_margin - right_margin, Inches(0.6))
    ctf = conclusion_box.text_frame
    ctf.text = '結論：建立模組化、可觀測且具回饋迴路的 Agent 架構，可加速迭代並降低風險。 如需進一步討論或範本，請聯絡。'
    for p in ctf.paragraphs:
        p.font.size = Pt(10)
        p.font.name = body_font
        p.font.color.rgb = slate_gray

    # Contact on same line (smaller)
    contact_box = slide.shapes.add_textbox(prs.slide_width - Inches(4.0), prs.slide_height - Inches(0.5), Inches(3.6), Inches(0.4))
    ctb = contact_box.text_frame
    ctb.text = 'contact@yourcompany.com'
    for p in ctb.paragraphs:
        p.font.size = Pt(10)
        p.font.name = body_font
        p.font.color.rgb = deep_blue

    # Save PPTX
    prs.save(pptx_path)

    # --- Create PNG representation using PIL ---
    # We'll create a 1280x720 image and draw similar content for preview
    W, H = 1280, 720
    img = Image.new('RGB', (W, H), color=(255,255,255))
    draw = ImageDraw.Draw(img)

    # Load fonts - fallback to default if not available
    try:
        font_title = ImageFont.truetype('arial.ttf', 36)
        font_sub = ImageFont.truetype('arial.ttf', 14)
        font_module = ImageFont.truetype('arial.ttf', 12)
        font_bold = ImageFont.truetype('arialbd.ttf', 14)
        font_small = ImageFont.truetype('arial.ttf', 11)
    except Exception:
        font_title = ImageFont.load_default()
        font_sub = ImageFont.load_default()
        font_module = ImageFont.load_default()
        font_bold = ImageFont.load_default()
        font_small = ImageFont.load_default()

    # Draw header
    draw.text((40,20), 'Agent 架構', font=font_title, fill=(28,40,51))
    draw.text((40,70), f'On-page Report — 簡要概覽    {date_str}    {author}', font=font_sub, fill=(80,90,100))

    # Diagram area
    dx, dy = 40, 120
    mod_w, mod_h = 150, 50
    gap = 25
    modules_short = ['感知', '理解/推理', '記憶', '學習', '執行']
    modules_desc = ['輸入資料解析', '策略決策與推理', '會話與知識庫', '模型更新/訓練', '執行/呼叫API']
    for i, m in enumerate(modules_short):
        x = dx + i*(mod_w+gap)
        y = dy
        rect_color = (236,243,249) if i%2==0 else (230,240,255)
        draw.rounded_rectangle([x,y,x+mod_w,y+mod_h], radius=10, fill=rect_color, outline=(46,64,83), width=2)
        # Title
        draw.text((x+10, y+8), m, font=font_bold, fill=(28,40,51))
        # small desc under
        draw.text((x+10, y+28), modules_desc[i], font=font_small, fill=(80,90,100))
        # arrows
        if i < len(modules_short)-1:
            ax1 = x+mod_w+5
            ay = y+mod_h/2
            ax2 = x+mod_w+gap-5
            draw.line([(ax1,ay),(ax2,ay)], fill=(46,64,83), width=3)
            # arrow head
            draw.polygon([(ax2,ay),(ax2-8,ay-6),(ax2-8,ay+6)], fill=(46,64,83))
    # Double arrow between memory(2) and learning(3)
    mem_x = dx + 2*(mod_w+gap)
    dbl_x1 = mem_x+mod_w+5
    dbl_x2 = dbl_x1 + (gap-10)
    draw.line([(dbl_x1, dy-20),(dbl_x2, dy-20)], fill=(94,168,167), width=3)
    draw.polygon([(dbl_x1,dy-20),(dbl_x1+8,dy-14),(dbl_x1+8,dy-26)], fill=(94,168,167))
    draw.polygon([(dbl_x2,dy-20),(dbl_x2-8,dy-14),(dbl_x2-8,dy-26)], fill=(94,168,167))

    # External boxes
    draw.rectangle([10,dy-30,140,dy+10], outline=(46,64,83), width=2)
    draw.text((18,dy-24), 'Users / Integrations', font=font_small, fill=(46,64,83))
    draw.text((18,dy-10), '(使用者 / 系統整合)', font=font_small, fill=(80,90,100))

    draw.rectangle([dx + 4*(mod_w+gap) + 130, dy+mod_h+10, dx + 4*(mod_w+gap) + 320, dy+mod_h+50], outline=(46,64,83), width=2)
    draw.text((dx + 4*(mod_w+gap) + 140, dy+mod_h+12), 'Environment', font=font_small, fill=(46,64,83))
    draw.text((dx + 4*(mod_w+gap) + 140, dy+mod_h+28), '(執行環境 / 感測場域)', font=font_small, fill=(80,90,100))

    # Right column bullets
    bx = 840
    by = 120
    draw.text((bx, by-20), '要點', font=font_bold, fill=(28,40,51))
    bullets = [
        ('模組職責明確','各模組單責任、易於測試與替換'),
        ('資料流設計','事件/訊息驅動，明確輸入/輸出契約'),
        ('可擴充性','模組化與 API 化，支援水平擴展'),
        ('安全與隱私','資料分類、存取控制與加密'),
        ('監控與回饋','日誌/指標/回饋迴路以支援改進'),
        ('部署考量','容器化、版本管理與 A/B 測試')
    ]
    ty = by
    for title,detail in bullets:
        draw.text((bx, ty), '• ' + title, font=font_bold, fill=(28,40,51))
        ty += 22
        draw.text((bx+12, ty), detail, font=font_small, fill=(80,90,100))
        ty += 34

    # Bottom conclusion and contact
    draw.line([(40, H-90),(W-40,H-90)], fill=(220,220,220), width=1)
    draw.text((40, H-70), '結論：建立模組化、可觀測且具回饋迴路的 Agent 架構，可加速迭代並降低風險。', font=font_small, fill=(80,90,100))
    draw.text((W-300, H-70), 'contact@yourcompany.com', font=font_small, fill=(28,40,51))

    img.save(png_path)

    print(f"Saved PPTX to {pptx_path}")
    print(f"Saved PNG to {png_path}")
    print('DONE')
except Exception as e:
    import traceback
    print('ERROR', repr(e))
    traceback.print_exc()