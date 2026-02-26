import os
try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception as e:
    print('Error: python-pptx not available:', e)
    raise

try:
    # Prepare directories
    os.makedirs('charts', exist_ok=True)
    os.makedirs('output', exist_ok=True)

    # Data
    years = list(range(2015, 2024))
    population = [23492074, 23539816, 23571227, 23588932, 23603121, 23561236, 23375314, 23264640, 23420442]

    # Try to create PNG chart via plotly; fallback to matplotlib
    png_path = 'charts/taiwan_population_plot.png'
    created_png = False
    try:
        import plotly.graph_objects as go
        import plotly.io as pio
        fig = go.Figure()
        fig.add_trace(go.Scatter(x=years, y=population, mode='lines+markers', line=dict(color='royalblue', width=3)))
        fig.update_layout(xaxis=dict(dtick=1), yaxis=dict(tickformat=','), template='plotly_white', width=900, height=506)
        # Attempt to write image
        pio.write_image(fig, png_path)
        created_png = True
        print('Created PNG via plotly:', png_path)
    except Exception as e:
        print('plotly image export failed, falling back to matplotlib:', e)
        try:
            import matplotlib.pyplot as plt
            plt.figure(figsize=(9,5))
            plt.plot(years, [p/1e6 for p in population], marker='o', color='royalblue')
            plt.xticks(years)
            plt.xlabel('Year')
            plt.ylabel('Population (millions)')
            plt.grid(True, linestyle='--', alpha=0.5)
            plt.tight_layout()
            plt.savefig(png_path, dpi=150)
            plt.close()
            created_png = True
            print('Created PNG via matplotlib:', png_path)
        except Exception as e2:
            print('matplotlib fallback failed:', e2)
            created_png = False

    # Create PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333333)  # 16:9
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "台灣近年人口趨勢分析（2015–2023）"
    subtitle.text = "自動產生 PPT - 包含圖表與數據摘要"

    # Slide 2: Key points
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '要點摘要'
    body = slide.shapes.placeholders[1].text_frame
    body.text = '1. 2015–2019：人口略增，接近高點。'
    p = body.add_paragraph()
    p.text = '2. 2020 起：出生率下降、天然增加為負，人口出現下滑。'
    p.level = 1
    p = body.add_paragraph()
    p.text = '3. 2023：淨遷入增加，導致人口回升。'

    # Slide 3: Data table
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = '年末總人口（2015–2023）'
    rows = 10
    cols = 2
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(6)
    height = Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(4)
    table.cell(0,0).text = '年份'
    table.cell(0,1).text = '年末總人口(人)'
    data = [
        (2015, 23492074),
        (2016, 23539816),
        (2017, 23571227),
        (2018, 23588932),
        (2019, 23603121),
        (2020, 23561236),
        (2021, 23375314),
        (2022, 23264640),
        (2023, 23420442),
    ]
    for i, (y, val) in enumerate(data, start=1):
        table.cell(i,0).text = str(y)
        table.cell(i,1).text = f"{val:,}"

    # Slide 4: Chart image
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = '人口趨勢圖（2015–2023）'
    if created_png and os.path.exists(png_path):
        pic_left = Inches(6)
        pic_top = Inches(1.5)
        slide.shapes.add_picture(png_path, pic_left, pic_top, width=Inches(6))
    else:
        # If no image, add small note
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(2))
        tf = tx_box.text_frame
        tf.text = '圖表檔案不存在，請確認 charts/taiwan_population_plot.png 是否可用。'

    # Slide 5: Interpretation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '解讀與影響'
    body = slide.shapes.placeholders[1].text_frame
    body.text = '1. 少子化趨勢明顯，出生數逐年降低。'
    p = body.add_paragraph(); p.text = '2. 老齡化與扶養比上升，長期社會支出壓力增加。'; p.level=1
    p = body.add_paragraph(); p.text = '3. 移民政策與經濟因素會影響短期人口波動。'; p.level=1

    # Slide 6: Conclusions & recommendations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '結論與建議'
    body = slide.shapes.placeholders[1].text_frame
    body.text = '1. 強化少子化相關政策，如托育、育兒津貼等。'
    body.add_paragraph(); body.paragraphs[1].text = '2. 吸引高技能移工與人才，緩解短期人口流失。'; body.paragraphs[1].level = 1
    body.add_paragraph(); body.paragraphs[2].text = '3. 調整長期社會保險與養老制度因應老年人口增加。'; body.paragraphs[2].level = 1

    # Slide 7: Sources
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = '資料來源'
    body = slide.shapes.placeholders[1].text_frame
    body.text = '內政部戶政司（Household Registration, Ministry of the Interior）'
    p = body.add_paragraph()
    p.text = 'Taiwan News, Taipei Times 與維基百科整理'
    p.level = 1

    pptx_path = 'output/taiwan_population_summary.pptx'
    prs.save(pptx_path)
    print('Saved PPTX to', pptx_path)
except Exception as ex:
    print('Error creating PPTX:', ex)

print('Done')