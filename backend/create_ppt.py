from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

try:
    os.makedirs('downloads', exist_ok=True)

    # Create presentation
    prs = Presentation()
    # Set slide size to 16:9 (in inches)
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # Slide layout blank (use first layout and then clear)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    # Title
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12.333333)
    title_height = Inches(1.0)
    title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    title_tf = title_box.text_frame
    title_tf.text = 'Agent 架構'
    title_par = title_tf.paragraphs[0]
    title_par.runs[0].font.size = Pt(32)
    title_par.runs[0].font.bold = True
    title_par.runs[0].font.name = 'Arial'

    # Left-side content box (text)
    left_left = Inches(0.5)
    left_top = Inches(1.5)
    left_width = Inches(6.0)
    left_height = Inches(5.5)
    left_box = slide.shapes.add_textbox(left_left, left_top, left_width, left_height)
    left_tf = left_box.text_frame
    left_tf.word_wrap = True

    # Short introduction (1-2 sentences)
    intro = left_tf.paragraphs[0]
    intro.text = '簡介：Agent 架構定義了多個協同模組，負責感知環境、規劃行動、執行任務、管理記憶、處理溝通與學習優化。'
    intro.font.size = Pt(12)
    intro.font.name = 'Arial'

    # Add modules as bullets
    modules = [
        ('Perception (感知)', ['蒐集環境與感測器資料', '分析輸入以形成狀態估計']),
        ('Planner (規劃)', ['根據目標與限制產生行動序列', '評估策略與替代方案']),
        ('Executor (執行)', ['將計劃轉為具體動作', '監控執行並回報狀態']),
        ('Memory (記憶/資料)', ['儲存歷史經驗與狀態資訊', '支援檢索與更新']),
        ('Communication (溝通/I/O)', ['處理外部輸入與輸出接口', '協調多方資料交換']),
        ('Learning (學習/優化)', ['從經驗改進策略與模型', '線上/離線調整參數'])
    ]

    # Add a blank paragraph to separate intro and list
    left_tf.add_paragraph()

    for mod_title, duties in modules:
        p = left_tf.add_paragraph()
        p.text = mod_title
        p.level = 0
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = 'Arial'
        # responsibilities as nested bullets
        for duty in duties:
            sp = left_tf.add_paragraph()
            sp.text = duty
            sp.level = 1
            sp.font.size = Pt(12)
            sp.font.name = 'Arial'

    # Right-side placeholder for image
    # Position and size to be noted: we'll choose a rectangle starting at x=7.2in, y=1.5in, width=5.6in, height=5.5in
    img_left = Inches(7.2)
    img_top = Inches(1.5)
    img_width = Inches(5.6)
    img_height = Inches(5.5)

    # Add gray rectangle as placeholder
    shape = slide.shapes.add_shape(
        1, img_left, img_top, img_width, img_height  # MSO_SHAPE.RECTANGLE == 1 in many versions
    )
    # Some python-pptx distributions require using MSO_AUTO_SHAPE_TYPE, but using 1 generally works.
    # Fill color gray
    try:
        from pptx.enum.dml import MSO_THEME_COLOR
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(200,200,200)
    except Exception:
        # best-effort; ignore if not supported
        pass

    # Add overlay text inside rectangle with filename
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = '占位圖: outputs/images/agent_architecture_image.png'
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = 1  # center

    # Add a small note textbox under the rectangle with explicit dimensions
    note_left = img_left
    note_top = img_top + img_height + Inches(0.15)
    note_width = img_width
    note_height = Inches(0.6)
    note_box = slide.shapes.add_textbox(note_left, note_top, note_width, note_height)
    note_tf = note_box.text_frame
    note_tf.text = f'圖像占位區域位置與尺寸 (近似)：left={img_left.inches:.2f}in, top={img_top.inches:.2f}in, width={img_width.inches:.2f}in, height={img_height.inches:.2f}in'
    note_tf.paragraphs[0].font.size = Pt(10)
    note_tf.paragraphs[0].font.italic = True
    note_tf.paragraphs[0].font.name = 'Arial'

    out_path = 'downloads/agent_architecture_draft.pptx'
    prs.save(out_path)

    print(f'Saved presentation to: {out_path}')
    print(f'DOWNLOAD: http://localhost:7777/download/agent_architecture_draft.pptx')
    print('Done')

except Exception as e:
    print('Error:', str(e))
    raise

print()